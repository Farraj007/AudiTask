import os,sys
from pptx import Presentation	
import datetime
import xlsxwriter 
from pptx.util import Cm
from tkinter import Tk 
from tkinter.filedialog import SaveFileDialog
# import win32com.client 
# from win32com.client import constants as c


now = datetime.datetime.now()
time_now=now.strftime("%d/%m/%Y %H:%M:%S")
#'Python_test_template.pptx'
prs = Presentation()
  
lines = {}
with open('01_Input/Config.txt') as f:
    for line in f:
        key, value = line.strip().split(':')
        lines[key] = (value)
name=lines.get('Name')
title=lines.get('Title')
manufacturer=lines.get('Manufacturer')
country=lines.get('Country')
city=lines.get('City')

first_slide_layout = prs.slide_layouts[0]
second_slide_layout= prs.slide_layouts[5]
third_slide_layout= prs.slide_layouts[5]
# s=prs.slides[0].shapes[0]
slide = prs.slides.add_slide(first_slide_layout)

top_part=slide.shapes.add_textbox( Cm(1.8), Cm(0.3), Cm(20), Cm(2))
note=top_part.text=f"{name} -{time_now}"

slide.placeholders[0].text =title

slide.placeholders[1].text = f"Name of author:{name}\nDate : {time_now}"
second_slide=prs.slides.add_slide(second_slide_layout)

top_part=second_slide.shapes.add_textbox( Cm(2), Cm(0.5), Cm(20), Cm(2))
note=top_part.text=f"{name} -{time_now}"


motors_folder= '01_Input'

motor_types_list,pictures_list=[],[]

for motor_types in os.listdir(motors_folder):
    motor_folder=f'01_Input/{motor_types}'

    if len(motor_folder.split('.'))==1:
        motor_types_list.append(motor_types)
        for picture in os.listdir(motor_folder+"/Pictures"):
            pictures_list.append(motor_folder+"/Pictures/"+picture)
            
        # for data_file in os.listdir(motor_folder+"/Data"):
        #     print(data_file)
motors_list=[]
count=0
for motor in pictures_list:    
    count+=1
    if count%2 == 0: 
        motors_list.append(motor.split('/')[3].split('_')[2])

second_slide.placeholders[0].text='Motor Types table'

# second_slide.shapes.add_table(len(motor_types_list)+1,4,Cm(0.5),Cm(4),Cm(20),10*len(motor_types_list))

workbook = xlsxwriter.Workbook('table_data.xlsx')
worksheet = workbook.add_worksheet()
table=worksheet.add_table(f'B2:E{len(motor_types_list)+2}',{'autofilter': False,
                                                            'style': 'Table Style Light 10',
                                                            'columns': [{'header': 'Motor Type'},
                                                                        {'header': 'Manufacturer'},
                                                                        {'header': 'Country'},
                                                                        {'header': 'City'},
                                                                       ],
                                            
                                                            })
header_format=workbook.add_format({'bold': True,
                                   'font_color': 'white',
                                   'bg_color':'#d45e6a',
                                   'align':'center',
                                   'valign':'center',
                                   'border':2,
                                   })
outer_cells_format=workbook.add_format({'align':'center',
                                        'valign':'vcenter',
                                        'border':2,
                                        })
middle_format=workbook.add_format({'bg_color':'#cad4e6',
                                   'align':'center',
                                   'valign':'vcenter',
                                   'border':2,
                                   })

worksheet.set_column('B2:B2', 10)
worksheet.set_column('C2:E2', 15)
worksheet.write_row('B2',['Motor Type','Manufacturer','Country','City'],header_format)
worksheet.merge_range(f'C3:C{len(motor_types_list)+2}',manufacturer, outer_cells_format)
worksheet.merge_range(f'D3:D{len(motor_types_list)+2}',country, middle_format)
worksheet.merge_range(f'E3:E{len(motor_types_list)+2}',city, outer_cells_format)

count=0
for motor in motors_list:
    count+=1
    if count%2 != 0:
        worksheet.write_row(f'B{count+2}',{motors_list[count-1]},middle_format)
    else:
        worksheet.write_row(f'B{count+2}',{motors_list[count-1]},outer_cells_format)   
workbook.close()
# i used this youtube as a refrence and understood what have done but wasn't done on my device cause of installation issue and i don't to waste time and uninstall everything since i installed anaconda here with the proffesor and anaconda did some stuff to my installation . never teh less here is the video refrence https://www.youtube.com/watch?v=_AiBCultl6U
#Grab the Active Instance of Excel.
# ExcelApp = win32com.client.GetActiveObject("Excel.Application")

# # Grab the workbook with the charts.
# xlWorkbook = ExcelApp.Workbooks("table_data.xlsx")
# for xlWorksheet in xlWorkbook.Worksheets:

#     # Grab the ChartObjects Collection for each sheet.
#     xlCharts = xlWorksheet.ChartObjects()
    
#     # Loop through each Chart in the ChartObjects Collection.
#     for index, xlChart in enumerate(xlCharts):

#         # Each chart needs to be on it's own slide, so at this point create a new slide.
#         PPTSlide = second_slide.Add(Index = index + 1, Layout = 12)  # 12 is a blank layout
        
#         # Display something to the user.
#         print('Exporting Chart {} from Worksheet {}'.format(xlChart.Name, xlWorksheet.Name))

#         # Copy the chart.
#         xlChart.Copy()

#         # Paste the Object to the Slide
#         PPTSlide.Shapes.PasteSpecial(DataType = 1) 
 
table=second_slide.shapes.add_table(len(motor_types_list)+1,4,Cm(2),Cm(5),Cm(24),10*len(motor_types_list)).table
header_values=['Motor Types ','Manufacturer','Country','city']
counter=0
for i in header_values:
    cell = table.cell(0, counter)
    print(cell)
    # cell.margin_left(Cm(0.5))
    # cell.margin_top(Cm(0.5))
    cell.text= i
    counter+=1
counter=0
for i in motor_types_list:
    counter+=1
    cell = table.cell(counter, 0)
    # cell.margin_left(Cm(0.5))
    # cell.margin_top(Cm(0.5))
    cell.text= i
    
for i in range(1,4):    
    cell = table.cell(1, i)
    other_cell = table.cell(len(motor_types_list),i )
    other_cell.merge(cell)
    cell.text= list(lines.values())[i+1]
    # cell.margin_left(Cm(0.5*len(motor_types_list)))
    # cell.margin_top(Cm(0.5*len(motor_types_list)))

# layout_swot = prs.slide_layouts[5]  # I use layout 'ID=5'
# for shape in layout_swot.placeholders:
#     print('%d %s' % (shape.placeholder_format.idx, shape.name))

# slide.shapes.add_chart_from_spreadsheet(table,
#     Cm(0.5), Cm(1.75), Cm(9), Cm(5))

for motor in motors_list:
    images_slide=prs.slides.add_slide(third_slide_layout)
    top_part=images_slide.shapes.add_textbox( Cm(1.8), Cm(0.3), Cm(20), Cm(2))
    note=top_part.text=f"{name} -{time_now}"
    images_slide.placeholders[0].text=motor
    images_slide.shapes.add_picture(pictures_list.pop(0),Cm(0), Cm(4), Cm(14),Cm(12))
    images_slide.shapes.add_picture(pictures_list.pop(0),Cm(18), Cm(4), Cm(14), Cm(12))


    
# Saving file in the same level with the same name
# prs.save('Python_test_template.pptx')



#taking the path from the user  
Tk().withdraw() 
filepath = SaveFileDialog()
prs.save(f'{filepath}/Python_test_template.pptx')
#taking the path from the user  
# from tkinter import *
# from tkinter import filedialog

# # Create an instance of window
# win=Tk()

# # Set the geometry of the window
# win.geometry("700x300")

# # Create a label
# Label(win, text="Click the button to open a dialog", font='Arial 16 bold').pack(pady=15)

# # Function to open a file in the system
# filepath = filedialog.SaveFileDialog(title="Choose where do you want to save ")
  
# # Create a button to trigger the dialog
# button = Button(win, text="Save", command=prs.save(f'{filepath}+Python_test_template.pptx'))
# button.pack()

# win.mainloop()


print("done")
